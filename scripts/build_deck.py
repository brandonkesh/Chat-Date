"""Builds the complete Crush product deck: previous features + today's new features + intro."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Brand system ----
NAVY   = RGBColor(0x1A, 0x22, 0x33)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
CARD   = RGBColor(0xF8, 0xFA, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)

HEAD = "Outfit"
BODY = "DM Sans"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    r._element.addprevious(r._element)  # keep as background
    return s


def _set_font(run, size, color, bold, font, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def text(s, l, t, w, h, runs, size=18, color=NAVY, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=4,
         italic=False):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, opt in para:
            r = p.add_run(); r.text = txt
            _set_font(r, opt.get("size", size), opt.get("color", color),
                      opt.get("bold", bold), opt.get("font", font), opt.get("italic", italic))
    return tb


def bar(s, l, t, w, h, color):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


def pill(s, l, t, w, h, label, fill, txtcolor=WHITE):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.adjustments[0] = 0.5
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.fill.background(); r.shadow.inherit = False
    tf = r.text_frame; tf.word_wrap = False
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    _set_font(run, 12, txtcolor, True, HEAD)
    return r


def card(s, l, t, w, h, fill=CARD):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    r.adjustments[0] = 0.06
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); r.line.width = Pt(1)
    r.shadow.inherit = False
    return r


def header(s, emoji_title, subtitle, new=False):
    bar(s, Inches(0), Inches(0), Inches(0.28), SH, BLUE)
    text(s, Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.9),
         emoji_title, size=34, color=NAVY, bold=True, font=HEAD)
    bar(s, Inches(0.75), Inches(1.42), Inches(1.3), Inches(0.06), ORANGE)
    if subtitle:
        text(s, Inches(0.7), Inches(1.55), Inches(11.6), Inches(0.6),
             subtitle, size=16, color=GRAY, font=BODY, italic=True)
    if new:
        pill(s, Inches(11.3), Inches(0.6), Inches(1.5), Inches(0.42), "★ NEW TODAY", ORANGE)


def feature_card(s, l, t, w, h, title, desc, accent=BLUE, new=False):
    card(s, l, t, w, h)
    pad = Inches(0.28)
    if new:
        pill(s, l + w - Inches(1.15), t + Inches(0.22), Inches(0.85), Inches(0.34), "NEW", ORANGE)
    text(s, l + pad, t + Inches(0.22), w - Inches(1.2), Inches(0.6),
         title, size=17, color=NAVY, bold=True, font=HEAD)
    text(s, l + pad, t + Inches(0.82), w - pad - Inches(0.15), h - Inches(0.95),
         desc, size=13.5, color=GRAY, font=BODY, line_spacing=1.05)


def group(s, l, t, w, title, items, accent=BLUE):
    text(s, l, t, w, Inches(0.5), title, size=18, color=accent, bold=True, font=HEAD)
    runs = [[("•  ", {"color": accent, "bold": True}), (it, {"color": NAVY})] for it in items]
    text(s, l, t + Inches(0.55), w, Inches(4.5), runs, size=14, color=NAVY,
         font=BODY, line_spacing=1.12, space_after=6)


# ============ SLIDE 1 — COVER ============
s = slide(NAVY)
bar(s, 0, 0, SW, SH, NAVY)
text(s, Inches(0), Inches(1.9), SW, Inches(1.2), "🔥", size=70, align=PP_ALIGN.CENTER, color=ORANGE)
text(s, Inches(0), Inches(3.05), SW, Inches(1.0),
     [[("Crush", {"size": 66, "bold": True, "color": WHITE, "font": HEAD})]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(4.2), SW, Inches(0.6),
     "Find your Crush without the wait.", size=24, color=RGBColor(0xCB,0xD5,0xE1),
     align=PP_ALIGN.CENTER, font=BODY)
text(s, Inches(0), Inches(4.85), SW, Inches(0.5),
     "The modern dating app that puts conversation first", size=16,
     color=GRAY, align=PP_ALIGN.CENTER, font=BODY)
pill(s, SW/2 - Inches(1.4), Inches(5.7), Inches(2.8), Inches(0.5), "crushmatchup.com", BLUE)

# ============ SLIDE 2 — INTRODUCTION (new) ============
s = slide()
header(s, "👋 Introduction", "What this overview covers")
text(s, Inches(0.75), Inches(2.3), Inches(11.6), Inches(1.4),
     [[("Crush is a complete, modern dating app — swipe to meet people, chat in real time, "
        "and get a little help from AI along the way. ", {"color": NAVY}),
       ("This deck walks through everything the app does today.", {"color": NAVY, "bold": True})]],
     size=18, color=NAVY, font=BODY, line_spacing=1.2)
c1 = Inches(0.75); c2 = Inches(6.95); cw = Inches(5.6); ct = Inches(3.5); ch = Inches(3.0)
card(s, c1, ct, cw, ch, LIGHT)
text(s, c1+Inches(0.35), ct+Inches(0.3), cw-Inches(0.6), Inches(0.6),
     "📘  Everything so far", size=19, color=BLUE, bold=True, font=HEAD)
group(s, c1+Inches(0.35), ct+Inches(1.0), cw-Inches(0.7),
      "", ["The full set of features already in the app",
           "Discovery, chat, AI tools, games & safety",
           "Membership plans and pricing"], accent=BLUE)
card(s, c2, ct, cw, ch, RGBColor(0xFF,0xF3,0xE8))
pill(s, c2+cw-Inches(1.5), ct+Inches(0.32), Inches(1.2), Inches(0.36), "★ NEW", ORANGE)
text(s, c2+Inches(0.35), ct+Inches(0.3), cw-Inches(1.7), Inches(0.6),
     "🆕  New this update", size=19, color=ORANGE, bold=True, font=HEAD)
group(s, c2+Inches(0.35), ct+Inches(1.0), cw-Inches(0.7),
      "", ["12 brand-new features shipped today",
           "New ways to meet, connect & have fun",
           "Clearly marked with a NEW tag throughout"], accent=ORANGE)

# ============ SLIDE 3 — WHAT IS CRUSH ============
s = slide()
header(s, "💙 What is Crush?", "A dating app built around real connection — not endless waiting.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "The Experience",
      ["Tinder-style swiping that's fast and fun",
       "Smart compatibility matching — interests, lifestyle, goals & location",
       "Real-time chat, voice notes, video calls & micro-dates",
       "A friendly, welcoming community"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "Why People Love It",
      ["30-day FREE trial of premium chat — no credit card needed",
       "Conversation-first: features that break the ice for you",
       "AI help whenever you want it",
       "Safety and trust built in from day one"], accent=ORANGE)

# ============ SLIDE 4 — CORE FEATURES ============
s = slide()
header(s, "✨ Core Features", "Everything you need to meet someone great.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "Discover",
      ["Swipe to like or pass", "Smart match scoring", "Top Picks & recommendations",
       "Advanced search filters", "Save profiles • hide profiles"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "Connect",
      ["Real-time messaging", "Voice notes & voice intros", "Intro videos on profiles",
       "Video calls (WebRTC)", "Micro-Dates — 5-minute virtual dates"], accent=ORANGE)

# ============ SLIDE 5 — AI TOOLS ============
s = slide()
header(s, "🤖 AI-Powered Dating Tools", "Your personal dating coach, built right in.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "Coaching & Advice",
      ["AI Dating Advisor — chat by voice or text for ideas & advice",
       "AI Conversation Coach — never run out of things to say",
       "AI Profile Optimizer — better bio & photo strategy"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "Smart & Safe",
      ["AI Scam Detector — real-time warnings on suspicious messages",
       "AI Photo Match — find compatible looks",
       "Personalized recommendations"], accent=ORANGE)

# ============ SLIDE 6 — GAMES & FUN ============
s = slide()
header(s, "🎮 Games & Fun", "Breaking the ice has never been easier.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "Chat Games",
      ["Two Truths and a Lie", "Emoji Story — a story in emojis only",
       "Would You Rather", "Date Bingo — a playful first-date checklist"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "Challenges & Quizzes",
      ["Personality Quiz — earn profile badges",
       "Daily login rewards & streaks",
       "Second Chance — revisit a pass"], accent=ORANGE)

# ============ SLIDE 7 — STAYING CONNECTED ============
s = slide()
header(s, "💞 Staying Connected", "Little nudges that keep the spark alive.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "Momentum",
      ["Daily login rewards — build a streak, earn Top Picks",
       "\"Your turn\" reminders when a match is waiting",
       "Likes You — see who's interested"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "After the Date",
      ["Date check-ins with a trusted-contact safety option",
       "Post-date feedback — rate how it went",
       "Success Stories — celebrate real couples"], accent=ORANGE)

# ============ SLIDE 8 — SAFETY & TRUST ============
s = slide()
header(s, "🛡️ Safety & Trust", "Dating should feel safe. We make sure it does.")
group(s, Inches(0.75), Inches(2.4), Inches(5.7), "Verified People",
      ["Photo verification with badge", "Age verification indicators",
       "VIP & trust badges", "Email verification"], accent=BLUE)
group(s, Inches(6.95), Inches(2.4), Inches(5.7), "Protected Chats",
      ["AI scam detection in real time", "Block & report anyone instantly",
       "Bidirectional blocking enforced", "Optional App Lock password"], accent=ORANGE)

# ============ SLIDE 9 — WHAT'S NEW (section intro) ============
s = slide(NAVY)
bar(s, 0, 0, SW, SH, NAVY)
pill(s, SW/2 - Inches(1.5), Inches(1.7), Inches(3.0), Inches(0.55), "★ SHIPPED TODAY", ORANGE)
text(s, Inches(0), Inches(2.6), SW, Inches(1.0),
     [[("What's New", {"size": 54, "bold": True, "color": WHITE, "font": HEAD})]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(3.8), SW, Inches(0.6),
     "12 brand-new features added in this update", size=20,
     color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER, font=BODY)
text(s, Inches(1.5), Inches(4.7), Inches(10.3), Inches(1.2),
     "New ways to meet, deeper ways to connect, and more reasons to come back every day — "
     "all built fresh today and shown on the next three slides.",
     size=15, color=GRAY, align=PP_ALIGN.CENTER, font=BODY, line_spacing=1.2)

# ---- helper for 2x2 new-feature grids ----
def new_grid(s, features):
    l1, l2 = Inches(0.75), Inches(6.95)
    t1, t2 = Inches(2.35), Inches(4.9)
    cw, ch = Inches(5.6), Inches(2.35)
    pos = [(l1, t1), (l2, t1), (l1, t2), (l2, t2)]
    for (title, desc), (l, t) in zip(features, pos):
        feature_card(s, l, t, cw, ch, title, desc, new=True)

# ============ SLIDE 10 — NEW: Connection & Personality ============
s = slide()
header(s, "🆕 New: Connection & Personality", "Deeper, more personal ways to click.", new=True)
new_grid(s, [
    ("💝 Kudos & \"Great Vibes\"",
     "Send a match some appreciation. Collect 3+ kudos and you earn a Great Vibes badge that shows you're a joy to talk to."),
    ("🎨 Dream Date Builder",
     "Pick the elements of your perfect date. Crush then shows what you and a match have in common — instant conversation starters."),
    ("🔮 Love Horoscopes",
     "A fresh, AI-written dating horoscope for your star sign every day. A fun little daily reason to open the app."),
    ("🎙️ Voice-First Mode",
     "Blur the photos in your feed so personality comes first. Get to know people by their words, not just their pictures."),
])

# ============ SLIDE 11 — NEW: Ways to Meet ============
s = slide()
header(s, "🆕 New: Ways to Meet", "Fresh, low-pressure ways to find someone.", new=True)
new_grid(s, [
    ("🎲 Blind Date Roulette",
     "Get instantly paired for a 5-minute surprise chat — photos hidden. At the end you both reveal; like each other and it's a match!"),
    ("❓ Question of the Week Club",
     "Everyone answers one fun prompt each week. Browse the answers and connect with people whose replies you love."),
    ("🐢 Slow Dating Mode",
     "Cap yourself at 5 likes a day to date more intentionally and give every match your full attention."),
    ("✅ Verified Badge in Chat",
     "See your match's verified status right inside the conversation, so you always know who you're talking to."),
])

# ============ SLIDE 12 — NEW: Community & Growth ============
s = slide()
header(s, "🆕 New: Community & Growth", "Keep the community buzzing and growing.", new=True)
new_grid(s, [
    ("🏆 Couple Leaderboard",
     "A playful weekly board of the chattiest couples on Crush (first names only) — a little friendly motivation to keep talking."),
    ("💞 Relationship Milestones",
     "Celebrate the moments that matter as a connection grows, with sweet in-app milestone moments."),
    ("🎁 Invite Links",
     "Share your own personal invite link with friends and watch the list of people who joined through you grow."),
    ("📊 Owner Dashboard",
     "A private admin view (owner only) with members, matches and activity at a glance — for running the app day to day."),
])

# ============ SLIDE 13 — MEMBERSHIP PLANS ============
s = slide()
header(s, "💎 Membership Plans", "Start with a 30-day FREE trial of premium chat — then pick your plan.")
plans = [
    ("Basic", "$4.99", "/mo", BLUE, ["Unlimited daily likes", "10 super likes / day",
        "See who viewed you", "Basic search filters", "Ad-free experience"]),
    ("Pro", "$9.99", "/mo", ORANGE, ["Everything in Basic", "AI Photo Match",
        "Advanced filters", "Priority in feed", "Read receipts"]),
    ("Elite", "$19.99", "/mo", GREEN, ["Everything in Pro", "Video calls",
        "Unlimited super likes", "Top of stack daily", "VIP badge & support"]),
]
lw, gap = Inches(3.9), Inches(0.28)
start = Inches(0.75)
top = Inches(2.35)
for i, (name, price, per, accent, items) in enumerate(plans):
    l = start + i * (lw + gap)
    card(s, l, top, lw, Inches(4.4))
    bar(s, l, top, lw, Inches(0.12), accent)
    text(s, l+Inches(0.3), top+Inches(0.35), lw-Inches(0.6), Inches(0.5),
         name, size=22, color=accent, bold=True, font=HEAD)
    text(s, l+Inches(0.3), top+Inches(0.95), lw-Inches(0.6), Inches(0.6),
         [[(price, {"size": 30, "bold": True, "color": NAVY, "font": HEAD}),
           (" "+per, {"size": 14, "color": GRAY})]])
    group(s, l+Inches(0.3), top+Inches(1.85), lw-Inches(0.55), "", items, accent=accent)

# ============ SLIDE 14 — CLOSING ============
s = slide(NAVY)
bar(s, 0, 0, SW, SH, NAVY)
text(s, Inches(0), Inches(1.9), SW, Inches(1.0), "🔥", size=64, align=PP_ALIGN.CENTER, color=ORANGE)
text(s, Inches(0), Inches(3.0), SW, Inches(1.0),
     [[("Ready to find your Crush?", {"size": 44, "bold": True, "color": WHITE, "font": HEAD})]],
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(4.1), SW, Inches(0.6),
     "Join free today — your first month of premium chat is on us.", size=18,
     color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER, font=BODY)
pill(s, SW/2 - Inches(1.5), Inches(5.0), Inches(3.0), Inches(0.55), "crushmatchup.com", BLUE)
text(s, Inches(0), Inches(5.8), SW, Inches(0.5),
     "No credit card required  •  Cancel anytime", size=13,
     color=GRAY, align=PP_ALIGN.CENTER, font=BODY)

import os
os.makedirs("exports", exist_ok=True)
out = "exports/Crush-Dating-App-Complete.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
